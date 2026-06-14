import collections 
import collections.abc
from pptx import Presentation

def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    with open("pptx_content.txt", "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"\n--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        f.write(text + "\n")

if __name__ == "__main__":
    extract_text(r"C:\Users\Thuy Quyen\Downloads\completeproduce\Tran_Thi_Thuy_Quyen_InternshipCCU_Week1.pdf.pptx")
